#!/usr/bin/env python3
"""Academic Style PPT - Day Care Medical Overview"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Colors
NAVY = RGBColor(15, 40, 80)
GOLD = RGBColor(196, 160, 112)
CREAM = RGBColor(250, 247, 240)
WHITE = RGBColor(255, 255, 255)
CHARCOAL = RGBColor(45, 45, 50)
GRAY = RGBColor(120, 120, 130)
LIGHT_GRAY = RGBColor(240, 240, 245)
BLUE = RGBColor(70, 130, 180)
GREEN = RGBColor(60, 140, 100)
RED = RGBColor(160, 50, 50)

def header(slide, title, sub=None):
    h = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.5))
    h.fill.solid(); h.fill.fore_color.rgb = NAVY; h.line.fill.background()
    g = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(1.5), Inches(13.333), Inches(0.06))
    g.fill.solid(); g.fill.fore_color.rgb = GOLD; g.line.fill.background()
    t = slide.shapes.add_textbox(Inches(0.8), Inches(0.35), Inches(11.5), Inches(0.7))
    p = t.text_frame.paragraphs[0]; p.text = title; p.font.size = Pt(28); p.font.bold = True; p.font.color.rgb = WHITE
    if sub:
        s = slide.shapes.add_textbox(Inches(0.8), Inches(1.0), Inches(11.5), Inches(0.4))
        p = s.text_frame.paragraphs[0]; p.text = sub; p.font.size = Pt(13); p.font.color.rgb = GOLD

def pg(slide, n): 
    p = slide.shapes.add_textbox(Inches(12.5), Inches(7.1), Inches(0.8), Inches(0.3))
    tf = p.text_frame.paragraphs[0]; tf.text = f"{n}/11"; tf.font.size = Pt(10); tf.font.color.rgb = GRAY; tf.alignment = PP_ALIGN.RIGHT

def bg(slide, color=CREAM):
    b = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, prs.slide_height)
    b.fill.solid(); b.fill.fore_color.rgb = color; b.line.fill.background()

# === SLIDE 1: Cover ===
s1 = prs.slides.add_slide(prs.slide_layouts[6])
bg(s1, NAVY)
l = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.15), prs.slide_height)
l.fill.solid(); l.fill.fore_color.rgb = GOLD; l.line.fill.background()
b = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(6.85), Inches(13.333), Inches(0.08))
b.fill.solid(); b.fill.fore_color.rgb = GOLD; b.line.fill.background()
c = s1.shapes.add_shape(MSO_SHAPE.OVAL, Inches(8.5), Inches(-1.5), Inches(7), Inches(7))
c.fill.solid(); c.fill.fore_color.rgb = RGBColor(25, 50, 90); c.line.fill.background()
badge = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.3), Inches(2.2), Inches(0.5))
badge.fill.solid(); badge.fill.fore_color.rgb = GOLD; badge.line.fill.background()
bt = s1.shapes.add_textbox(Inches(0.8), Inches(1.35), Inches(2.2), Inches(0.5))
p = bt.text_frame.paragraphs[0]; p.text = "学术报告"; p.font.size = Pt(14); p.font.bold = True; p.font.color.rgb = NAVY; p.alignment = PP_ALIGN.CENTER
mt = s1.shapes.add_textbox(Inches(0.8), Inches(2.3), Inches(10), Inches(1.4))
p = mt.text_frame.paragraphs[0]; p.text = "日间医疗"; p.font.size = Pt(72); p.font.bold = True; p.font.color.rgb = WHITE
st = s1.shapes.add_textbox(Inches(0.8), Inches(3.6), Inches(10), Inches(0.8))
p = st.text_frame.paragraphs[0]; p.text = "起源·发展·现状·未来趋势"; p.font.size = Pt(28); p.font.color.rgb = GOLD
ln = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(4.5), Inches(3), Inches(0.04))
ln.fill.solid(); ln.fill.fore_color.rgb = GOLD; ln.line.fill.background()
dt = s1.shapes.add_textbox(Inches(0.8), Inches(4.8), Inches(8), Inches(0.6))
p = dt.text_frame.paragraphs[0]; p.text = "基于国内外文献与政策文件的综合分析"; p.font.size = Pt(16); p.font.color.rgb = RGBColor(180, 180, 190)
d2 = s1.shapes.add_textbox(Inches(0.8), Inches(6.3), Inches(5), Inches(0.4))
p = d2.text_frame.paragraphs[0]; p.text = "2026年3月"; p.font.size = Pt(14); p.font.color.rgb = RGBColor(150, 150, 160)

# === SLIDE 2: TOC ===
s2 = prs.slides.add_slide(prs.slide_layouts[6]); bg(s2)
header(s2, "目 录", "CONTENTS"); pg(s2, 2)
toc = [("01","起源与概念界定","Origin and Concept"),("02","国际发展背景","International Context"),
       ("03","中国发展历程","Development History"),("04","发展现状与规模","Current Status"),
       ("05","主要运作模式","Operational Models"),("06","挑战与对策","Challenges"),
       ("07","未来发展趋势","Future Trends")]
y = 2.0
for num, cn, en in toc:
    nb = s2.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(y), Inches(0.8), Inches(0.6))
    nb.fill.solid(); nb.fill.fore_color.rgb = NAVY; nb.line.fill.background()
    nt = s2.shapes.add_textbox(Inches(0.8), Inches(y+0.1), Inches(0.8), Inches(0.5))
    p = nt.text_frame.paragraphs[0]; p.text = num; p.font.size = Pt(18); p.font.bold = True; p.font.color.rgb = GOLD; p.alignment = PP_ALIGN.CENTER
    ct = s2.shapes.add_textbox(Inches(1.8), Inches(y+0.05), Inches(5), Inches(0.4))
    p = ct.text_frame.paragraphs[0]; p.text = cn; p.font.size = Pt(18); p.font.bold = True; p.font.color.rgb = CHARCOAL
    et = s2.shapes.add_textbox(Inches(1.8), Inches(y+0.4), Inches(5), Inches(0.3))
    p = et.text_frame.paragraphs[0]; p.text = en; p.font.size = Pt(11); p.font.color.rgb = GRAY
    y += 0.7
db = s2.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(10), Inches(1.8), Inches(3), Inches(5.2))
db.fill.solid(); db.fill.fore_color.rgb = NAVY; db.line.fill.background()
gl = s2.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(10.1), Inches(2), Inches(0.06), Inches(4.8))
gl.fill.solid(); gl.fill.fore_color.rgb = GOLD; gl.line.fill.background()
qt = s2.shapes.add_textbox(Inches(10.4), Inches(2.5), Inches(2.4), Inches(3))
tf = qt.text_frame; tf.word_wrap = True; p = tf.paragraphs[0]
p.text = '"日间医疗是提高医疗效率、优化资源配置的重要途径"'; p.font.size = Pt(14); p.font.color.rgb = WHITE; p.font.italic = True

# === SLIDE 3: Origin ===
s3 = prs.slides.add_slide(prs.slide_layouts[6]); bg(s3)
header(s3, "一、起源与概念界定", "Origin and Concept Definition"); pg(s3, 3)
db = s3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.9), Inches(12.3), Inches(1.7))
db.fill.solid(); db.fill.fore_color.rgb = WHITE; db.line.color.rgb = GOLD; db.line.width = Pt(2)
ga = s3.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.9), Inches(0.12), Inches(1.7))
ga.fill.solid(); ga.fill.fore_color.rgb = GOLD; ga.line.fill.background()
dt = s3.shapes.add_textbox(Inches(0.9), Inches(2.0), Inches(11.5), Inches(0.4))
p = dt.text_frame.paragraphs[0]; p.text = "【日间医疗定义】Day Care / Ambulatory Care"; p.font.size = Pt(16); p.font.bold = True; p.font.color.rgb = NAVY
dx = s3.shapes.add_textbox(Inches(0.9), Inches(2.45), Inches(11.5), Inches(1))
tf = dx.text_frame; tf.word_wrap = True; p = tf.paragraphs[0]
p.text = "指患者在院时间不超过24小时的医疗服务和住院模式。患者当日入院、完成手术或操作、术后观察恢复，当日或次日出院的的一种高效医疗服务模式。" ; p.font.size = Pt(15); p.font.color.rgb = CHARCOAL
feats = [("⏱","住院时间短","≤24小时"),("💰","医疗费用低","降低成本"),("🏥","效率提升","床位周转快"),("👨‍⚕️","流程标准化","规范管理")]
x = 0.7
for icon, title, desc in feats:
    bx = s3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(3.9), Inches(2.9), Inches(1.6))
    bx.fill.solid(); bx.fill.fore_color.rgb = NAVY; bx.line.fill.background()
    it = s3.shapes.add_textbox(Inches(x), Inches(4.0), Inches(2.9), Inches(0.6))
    p = it.text_frame.paragraphs[0]; p.text = icon; p.font.size = Pt(24); p.alignment = PP_ALIGN.CENTER
    tt = s3.shapes.add_textbox(Inches(x), Inches(4.55), Inches(2.9), Inches(0.4))
    p = tt.text_frame.paragraphs[0]; p.text = title; p.font.size = Pt(15); p.font.bold = True; p.font.color.rgb = WHITE; p.alignment = PP_ALIGN.CENTER
    ds = s3.shapes.add_textbox(Inches(x), Inches(4.95), Inches(2.9), Inches(0.4))
    p = ds.text_frame.paragraphs[0]; p.text = desc; p.font.size = Pt(12); p.font.color.rgb = GOLD; p.alignment = PP_ALIGN.CENTER
    x += 3.2
ht = s3.shapes.add_textbox(Inches(0.5), Inches(5.8), Inches(12.3), Inches(0.8))
tf = ht.text_frame; tf.word_wrap = True; p = tf.paragraphs[0]
p.text = "📖 历史渊源：1909年苏格兰医生James Nicoll首次提出日间手术概念 | 1995年国际日间手术协会（IAASS）成立 | 2012年中国发布首个日间手术管理规范"; p.font.size = Pt(11); p.font.color.rgb = GRAY

# === SLIDE 4: International ===
s4 = prs.slides.add_slide(prs.slide_layouts[6]); bg(s4)
header(s4, "一、国际发展背景", "International Development Context"); pg(s4, 4)
countries = [("🇬🇧 英国","日间手术发源地",["1909年概念提出","90年代全面推广","日间手术占比超70%"],BLUE),
             ("🇺🇸 美国","全球领先",["ASC体系成熟","日间手术占比85%+","覆盖大部分术式"],GREEN),
             ("🇩🇪 德国","标准化典范",["1980年代开始发展","严格日间手术标准","高效管理系统"],RGBColor(140,100,180)),
             ("🇨🇳 中国","快速发展",["2012年规范化","2020年扩大试点","占比20-30%"],RED)]
x = 0.5
for country, tag, items, color in countries:
    cd = s4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(1.9), Inches(3), Inches(4.9))
    cd.fill.solid(); cd.fill.fore_color.rgb = WHITE; cd.line.color.rgb = color; cd.line.width = Pt(2)
    hd = s4.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(1.9), Inches(3), Inches(0.8))
    hd.fill.solid(); hd.fill.fore_color.rgb = color; hd.line.fill.background()
    nm = s4.shapes.add_textbox(Inches(x), Inches(2.0), Inches(3), Inches(0.6))
    p = nm.text_frame.paragraphs[0]; p.text = country; p.font.size = Pt(18); p.font.bold = True; p.font.color.rgb = WHITE; p.alignment = PP_ALIGN.CENTER
    tg = s4.shapes.add_textbox(Inches(x+0.2), Inches(2.9), Inches(2.6), Inches(0.4))
    p = tg.text_frame.paragraphs[0]; p.text = "◆ " + tag; p.font.size = Pt(13); p.font.color.rgb = color
    y = 3.4
    for item in items:
        it = s4.shapes.add_textbox(Inches(x+0.2), Inches(y), Inches(2.6), Inches(0.4))
        p = it.text_frame.paragraphs[0]; p.text = "• " + item; p.font.size = Pt(12); p.font.color.rgb = CHARCOAL
        y += 0.5
    x += 3.3

# === SLIDE 5: China History ===
s5 = prs.slides.add_slide(prs.slide_layouts[6]); bg(s5)
header(s5, "二、中国发展历程", "Development History in China"); pg(s5, 5)
timeline = [("2001-2010","萌芽期","• 少数大型三甲医院探索 | • 主要集中在眼科、日间化疗 | • 缺乏统一标准",NAVY),
            ("2012","规范期","• 原卫生部发布首个管理文件 | • 北京、上海开始试点 | • 成立日间手术合作联盟",BLUE),
            ("2015-2018","推广期","• 国家卫计委扩大试点 | • 多省市出台政策 | • 医保支付改革推进",GREEN),
            ("2019-至今","快车道","• 新技术推动（微创、麻醉） | • 疫情加速非住院化 | • 日间病房标准化建设",GOLD)]
ln = s5.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.5), Inches(3.15), Inches(10.3), Inches(0.03))
ln.fill.solid(); ln.fill.fore_color.rgb = NAVY; ln.line.fill.background()
for i, (period, stage, content, color) in enumerate(timeline):
    dot = s5.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.35), Inches(2.98), Inches(0.35), Inches(0.35))
    dot.fill.solid(); dot.fill.fore_color.rgb = color; dot.line.fill.background()
    y_pos = 2.2 if i % 2 == 0 else 3.5
    pb = s5.shapes.add_textbox(Inches(0.8), Inches(y_pos), Inches(1.5), Inches(0.4))
    p = pb.text_frame.paragraphs[0]; p.text = period; p.font.size = Pt(12); p.font.bold = True; p.font.color.rgb = color; p.alignment = PP_ALIGN.CENTER
    sb = s5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2), Inches(y_pos+0.4), Inches(1.8), Inches(0.5))
    sb.fill.solid(); sb.fill.fore_color.rgb = color; sb.line.fill.background()
    st = s5.shapes.add_textbox(Inches(2), Inches(y_pos+0.45), Inches(1.8), Inches(0.4))
    p = st.text_frame.paragraphs[0]; p.text = stage; p.font.size = Pt(13); p.font.bold = True; p.font.color.rgb = WHITE if color != GOLD else NAVY; p.alignment = PP_ALIGN.CENTER
    cb = s5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4), Inches(y_pos+0.4), Inches(8.5), Inches(1.1))
    cb.fill.solid(); cb.fill.fore_color.rgb = WHITE; cb.line.color.rgb = color; cb.line.width = Pt(1)
    ct = s5.shapes.add_textbox(Inches(4.2), Inches(y_pos+0.5), Inches(8.1), Inches(0.9))
    tf = ct.text_frame; tf.word_wrap = True; p = tf.paragraphs[0]; p.text = content; p.font.size = Pt(11); p.font.color.rgb = CHARCOAL

# === SLIDE 6: Current Status ===
s6 = prs.slides.add_slide(prs.slide_layouts[6]); bg(s6)
header(s6, "三、发展现状与规模", "Current Status and Scale"); pg(s6, 6)
stats = [("30%+","年均增长率","过去5年日间手术量快速增长",BLUE),
         ("20-30%","择期手术占比","与欧美国家差距明显，提升空间大",GREEN),
         (">10,000","开展机构数","全国三级医院广泛开展日间服务",NAVY),
         ("200+","术式种类","覆盖普外、眼科、骨科等多个专科",GOLD)]
y = 1.9
for num, label, desc, color in stats:
    cd = s6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y), Inches(6), Inches(1.2))
    cd.fill.solid(); cd.fill.fore_color.rgb = WHITE; cd.line.color.rgb = color; cd.line.width = Pt(2)
    bd = s6.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(y), Inches(0.12), Inches(1.2))
    bd.fill.solid(); bd.fill.fore_color.rgb = color; bd.line.fill.background()
    nm = s6.shapes.add_textbox(Inches(0.8), Inches(y+0.12), Inches(2), Inches(0.6))
    p = nm.text_frame.paragraphs[0]; p.text = num; p.font.size = Pt(30); p.font.bold = True; p.font.color.rgb = color
    lb = s6.shapes.add_textbox(Inches(0.8), Inches(y+0.7), Inches(2), Inches(0.4))
    p = lb.text_frame.paragraphs[0]; p.text = label; p.font.size = Pt(13); p.font.bold = True; p.font.color.rgb = CHARCOAL
    dc = s6.shapes.add_textbox(Inches(2.9), Inches(y+0.3), Inches(3.4), Inches(0.8))
    tf = dc.text_frame; tf.word_wrap = True; p = tf.paragraphs[0]; p.text = desc; p.font.size = Pt(12); p.font.color.rgb = GRAY
    y += 1.3
sp = s6.shapes.add_textbox(Inches(7), Inches(1.9), Inches(5.8), Inches(0.5))
p = sp.text_frame.paragraphs[0]; p.text = "📋 覆盖主要专科"; p.font.size = Pt(16); p.font.bold = True; p.font.color.rgb = NAVY
specialties = [("🏥","普外科","胆囊、疝气、阑尾"),("👁","眼科","白内障手术"),("🦴","骨科","关节镜、日间骨折"),
               ("👃","耳鼻喉科","腺样体、鼻窦手术"),("💊","肿瘤科","日间化疗"),("🫁","呼吸科","支气管镜"),
               ("💉","介入治疗","血管、穿刺"),("🧠","神经外科","介入手术")]
y = 2.4
for icon, dept, examples in specialties:
    rw = s6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7), Inches(y), Inches(5.8), Inches(0.55))
    rw.fill.solid(); rw.fill.fore_color.rgb = WHITE if int(y*10)%20 < 10 else LIGHT_GRAY; rw.line.fill.background()
    lb = s6.shapes.add_textbox(Inches(7.2), Inches(y+0.1), Inches(0.8), Inches(0.4))
    p = lb.text_frame.paragraphs[0]; p.text = icon; p.font.size = Pt(14)
    dt = s6.shapes.add_textbox(Inches(7.8), Inches(y+0.1), Inches(1.2), Inches(0.4))
    p = dt.text_frame.paragraphs[0]; p.text = dept; p.font.size = Pt(13); p.font.bold = True; p.font.color.rgb = CHARCOAL
    ep = s6.shapes.add_textbox(Inches(9.0), Inches(y+0.1), Inches(3.6), Inches(0.4))
    p = ep.text_frame.paragraphs[0]; p.text = examples; p.font.size = Pt(11); p.font.color.rgb = GRAY
    y += 0.58

# === SLIDE 7: Models ===
s7 = prs.slides.add_slide(prs.slide_layouts[6]); bg(s7)
header(s7, "四、主要运作模式", "Operational Models"); pg(s7, 7)
models = [("模式一","独立日间手术中心",["独立的日间手术中心","专用床位和手术室","专职医护团队","完全按日间模式运行"],NAVY),
          ("模式二","集中管理式",["医院统一管理日间患者","分散收治、统一排程","共享手术室资源","多学科协作(MDT)"],BLUE),
          ("模式三","科室嵌入式",["各临床科室设置日间床位","科室内独立运作","便于术后观察","适合复杂病例"],GREEN)]
x = 0.5
for mnum, mname, items, color in models:
    cd = s7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(1.9), Inches(4), Inches(5.1))
    cd.fill.solid(); cd.fill.fore_color.rgb = WHITE; cd.line.color.rgb = color; cd.line.width = Pt(2)
    hd = s7.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(1.9), Inches(4), Inches(1.2))
    hd.fill.solid(); hd.fill.fore_color.rgb = color; hd.line.fill.background()
    nm = s7.shapes.add_textbox(Inches(x), Inches(2.0), Inches(4), Inches(0.5))
    p = nm.text_frame.paragraphs[0]; p.text = mnum; p.font.size = Pt(14); p.font.color.rgb = GOLD; p.alignment = PP_ALIGN.CENTER
    nt = s7.shapes.add_textbox(Inches(x), Inches(2.45), Inches(4), Inches(0.6))
    p = nt.text_frame.paragraphs[0]; p.text = mname; p.font.size = Pt(18); p.font.bold = True; p.font.color.rgb = WHITE; p.alignment = PP_ALIGN.CENTER
    y = 3.3
    for item in items:
        bl = s7.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x+0.3), Inches(y+0.12), Inches(0.12), Inches(0.12))
        bl.fill.solid(); bl.fill.fore_color.rgb = color; bl.line.fill.background()
        it = s7.shapes.add_textbox(Inches(x+0.55), Inches(y), Inches(3.3), Inches(0.5))
        tf = it.text_frame; tf.word_wrap = True; p = tf.paragraphs[0]; p.text = item; p.font.size = Pt(13); p.font.color.rgb = CHARCOAL
        y += 0.65
    x += 4.3

# === SLIDE 8: Challenges ===
s8 = prs.slides.add_slide(prs.slide_layouts[6]); bg(s8)
header(s8, "四、挑战与对策", "Challenges and Countermeasures"); pg(s8, 8)
challenges = [("⚠️","患者认知不足","传统观念影响，对日间手术安全性存在担忧，术后家庭护理顾虑","加强科普宣传，建立信任",RED),
              ("💳","支付方式限制","医保支付政策不完善，日间手术定价不合理，报销比例待提高","推动支付方式改革",BLUE),
              ("🏥","质量安全保障","术后并发症风险，应急处理能力，随访管理难度","完善质量保障体系",GREEN),
              ("🏙","资源配置不均","城乡差距明显，不同地区发展不平衡，基层能力不足","促进优质资源下沉",GOLD)]
y = 1.9
for icon, title, problem, solution, color in challenges:
    cd = s8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y), Inches(12.3), Inches(1.25))
    cd.fill.solid(); cd.fill.fore_color.rgb = WHITE; cd.line.color.rgb = color; cd.line.width = Pt(1.5)
    ib = s8.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(y), Inches(0.8), Inches(1.25))
    ib.fill.solid(); ib.fill.fore_color.rgb = color; ib.line.fill.background()
    ic = s8.shapes.add_textbox(Inches(0.5), Inches(y+0.35), Inches(0.8), Inches(0.6))
    p = ic.text_frame.paragraphs[0]; p.text = icon; p.font.size = Pt(22); p.alignment = PP_ALIGN.CENTER
    tt = s8.shapes.add_textbox(Inches(1.5), Inches(y+0.15), Inches(2.5), Inches(0.4))
    p = tt.text_frame.paragraphs[0]; p.text = title; p.font.size = Pt(15); p.font.bold = True; p.font.color.rgb = CHARCOAL
    pr = s8.shapes.add_textbox(Inches(1.5), Inches(y+0.55), Inches(6), Inches(0.6))
    tf = pr.text_frame; tf.word_wrap = True; p = tf.paragraphs[0]; p.text = problem; p.font.size = Pt(11); p.font.color.rgb = GRAY
    ar = s8.shapes.add_textbox(Inches(7.8), Inches(y+0.35), Inches(0.4), Inches(0.5))
    p = ar.text_frame.paragraphs[0]; p.text = "→"; p.font.size = Pt(20); p.font.color.rgb = color
    sb = s8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.3), Inches(y+0.25), Inches(4.3), Inches(0.75))
    sb.fill.solid(); sb.fill.fore_color.rgb = color; sb.line.fill.background()
    sl = s8.shapes.add_textbox(Inches(8.3), Inches(y+0.4), Inches(4.3), Inches(0.5))
    p = sl.text_frame.paragraphs[0]; p.text = solution; p.font.size = Pt(13); p.font.bold = True; p.font.color.rgb = WHITE; p.alignment = PP_ALIGN.CENTER
    y += 1.35

# === SLIDE 9: Future Trends ===
s9 = prs.slides.add_slide(prs.slide_layouts[6]); bg(s9)
header(s9, "五、未来发展趋势", "Future Development Trends"); pg(s9, 9)
trends = [("🤖 智慧化","技术驱动效率提升","• 远程术后随访系统\n• AI辅助患者筛选\n• 智能预警与干预",BLUE),
          ("🏠 社区化","分级诊疗深度融合","• 社区医疗机构参与\n• 家庭病床+日间模式\n• 上下转诊无缝衔接",GREEN),
          ("💊 药学服务","安全用药保障","• 临床药师参与用药管理\n• 个体化用药方案\n• 药物安全性监测",RGBColor(140,100,180)),
          ("📊 标准化","质量持续改进","• 统一质量评价体系\n• 临床路径标准化\n• 数据上报与公示",GOLD)]
x = 0.5
for title, benefit, content, color in trends:
    cd = s9.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(1.9), Inches(6), Inches(4.9))
    cd.fill.solid(); cd.fill.fore_color.rgb = WHITE; cd.line.color.rgb = color; cd.line.width = Pt(2)
    hd = s9.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(1.9), Inches(6), Inches(1))
    hd.fill.solid(); hd.fill.fore_color.rgb = color; hd.line.fill.background()
    nm = s9.shapes.add_textbox(Inches(x), Inches(2.0), Inches(6), Inches(0.5))
    p = nm.text_frame.paragraphs[0]; p.text = title; p.font.size = Pt(18); p.font.bold = True; p.font.color.rgb = WHITE; p.alignment = PP_ALIGN.CENTER
    bf = s9.shapes.add_textbox(Inches(x), Inches(2.45), Inches(6), Inches(0.4))
    p = bf.text_frame.paragraphs[0]; p.text = benefit; p.font.size = Pt(12); p.font.color.rgb = GOLD; p.alignment = PP_ALIGN.CENTER
    ct = s9.shapes.add_textbox(Inches(x+0.3), Inches(3.1), Inches(5.4), Inches(3.5))
    tf = ct.text_frame; tf.word_wrap = True; p = tf.paragraphs[0]; p.text = content; p.font.size = Pt(13); p.font.color.rgb = CHARCOAL
    x += 6.3

# === SLIDE 10: Policy ===
s10 = prs.slides.add_slide(prs.slide_layouts[6]); bg(s10)
header(s10, "五、政策导向", "Policy Direction"); pg(s10, 10)
policies = [("📋 《日间手术管理规范》","国家卫健委发布的日间手术管理指导文件，明确准入标准、操作流程、质量控制，要求建立日间手术质量评价体系"),
            ("💳 支付方式改革","推进DRG/DIP付费与日间手术结合，探索按病种付费的日间手术定价机制，提高日间手术医保报销比例"),
            ("🏗 基础设施建设","推动日间手术中心标准化建设，鼓励三级医院建立独立日间手术部，支持日间手术信息化建设"),
            ("👥 人才队伍建设","加强日间手术专业人才培训，建立日间手术专科护士培训体系，推动多学科协作团队建设")]
y = 1.9
for title, content in policies:
    cd = s10.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y), Inches(12.3), Inches(1.25))
    cd.fill.solid(); cd.fill.fore_color.rgb = WHITE; cd.line.color.rgb = NAVY; cd.line.width = Pt(1.5)
    nm = s10.shapes.add_textbox(Inches(0.8), Inches(y+0.15), Inches(11.7), Inches(0.4))
    p = nm.text_frame.paragraphs[0]; p.text = title; p.font.size = Pt(15); p.font.bold = True; p.font.color.rgb = NAVY
    ct = s10.shapes.add_textbox(Inches(0.8), Inches(y+0.55), Inches(11.7), Inches(0.6))
    tf = ct.text_frame; tf.word_wrap = True; p = tf.paragraphs[0]; p.text = content; p.font.size = Pt(12); p.font.color.rgb = CHARCOAL
    y += 1.35

# === SLIDE 11: Summary ===
s11 = prs.slides.add_slide(prs.slide_layouts[6])
bg(s11, NAVY)
l = s11.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.15), prs.slide_height)
l.fill.solid(); l.fill.fore_color.rgb = GOLD; l.line.fill.background()
b = s11.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(6.85), Inches(13.333), Inches(0.08