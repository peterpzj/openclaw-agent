#!/usr/bin/env python3
"""
Academic Style PPT Generator for Day Care Medical Research
Enhanced design with professional academic aesthetics
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ============================================================================
# DESIGN SYSTEM - Academic Navy & Gold Theme
# ============================================================================

NAVY = RGBColor(15, 40, 80)
GOLD = RGBColor(196, 160, 112)
LIGHT_NAVY = RGBColor(32, 64, 128)
CREAM = RGBColor(250, 247, 240)
WHITE = RGBColor(255, 255, 255)
CHARCOAL = RGBColor(45, 45, 50)
GRAY = RGBColor(120, 120, 130)
LIGHT_GRAY = RGBColor(240, 240, 245)
ACCENT_BLUE = RGBColor(70, 130, 180)
ACCENT_GREEN = RGBColor(60, 140, 100)
ACCENT_RED = RGBColor(160, 50, 50)

def add_page_number(slide, num, total):
    pn = slide.shapes.add_textbox(Inches(12.5), Inches(7.1), Inches(0.8), Inches(0.3))
    tf = pn.text_frame
    p = tf.paragraphs[0]
    p.text = f"{num}/{total}"
    p.font.size = Pt(10)
    p.font.color.rgb = GRAY
    p.alignment = PP_ALIGN.RIGHT

def add_header_accent(slide, title_text, subtitle=None):
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.5))
    header.fill.solid()
    header.fill.fore_color.rgb = NAVY
    header.line.fill.background()
    
    gold_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(1.5), Inches(13.333), Inches(0.06))
    gold_bar.fill.solid()
    gold_bar.fill.fore_color.rgb = GOLD
    gold_bar.line.fill.background()
    
    title = slide.shapes.add_textbox(Inches(0.8), Inches(0.35), Inches(11.5), Inches(0.7))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    if subtitle:
        sub = slide.shapes.add_textbox(Inches(0.8), Inches(1.0), Inches(11.5), Inches(0.4))
        tf = sub.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(13)
        p.font.color.rgb = GOLD

def add_accent_line(slide, x, y, width, color=GOLD):
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(width), Inches(0.04))
    line.fill.solid()
    line.fill.fore_color.rgb = color
    line.line.fill.background()

# ============================================================================
# SLIDE 1: Cover
# ============================================================================
slide1 = prs.slides.add_slide(prs.slide_layouts[6])

bg = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, prs.slide_height)
bg.fill.solid()
bg.fill.fore_color.rgb = NAVY
bg.line.fill.background()

# Left gold bar
deco1 = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.15), prs.slide_height)
deco1.fill.solid()
deco1.fill.fore_color.rgb = GOLD
deco1.line.fill.background()

# Bottom gold line
deco2 = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(6.85), Inches(13.333), Inches(0.08))
deco2.fill.solid()
deco2.fill.fore_color.rgb = GOLD
deco2.line.fill.background()

# Decorative circle
circle = slide1.shapes.add_shape(MSO_SHAPE.OVAL, Inches(8.5), Inches(-1.5), Inches(7), Inches(7))
circle.fill.solid()
circle.fill.fore_color.rgb = RGBColor(25, 50, 90)
circle.line.fill.background()

# Badge
badge = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.3), Inches(2.2), Inches(0.5))
badge.fill.solid()
badge.fill.fore_color.rgb = GOLD
badge.line.fill.background()

badge_text = slide1.shapes.add_textbox(Inches(0.8), Inches(1.35), Inches(2.2), Inches(0.5))
tf = badge_text.text_frame
p = tf.paragraphs[0]
p.text = "学术报告"
p.font.size = Pt(14)
p.font.bold = True
p.font.color.rgb = NAVY
p.alignment = PP_ALIGN.CENTER

# Main title
main_title = slide1.shapes.add_textbox(Inches(0.8), Inches(2.3), Inches(10), Inches(1.4))
tf = main_title.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "日间医疗"
p.font.size = Pt(72)
p.font.bold = True
p.font.color.rgb = WHITE

# Subtitle
sub = slide1.shapes.add_textbox(Inches(0.8), Inches(3.6), Inches(10), Inches(0.8))
tf = sub.text_frame
p = tf.paragraphs[0]
p.text = "起源·发展·现状·未来趋势"
p.font.size = Pt(28)
p.font.color.rgb = GOLD

add_accent_line(slide1, 0.8, 4.5, 3)

desc = slide1.shapes.add_textbox(Inches(0.8), Inches(4.8), Inches(8), Inches(0.6))
tf = desc.text_frame
p = tf.paragraphs[0]
p.text = "基于国内外文献与政策文件的综合分析"
p.font.size = Pt(16)
p.font.color.rgb = RGBColor(180, 180, 190)

date = slide1.shapes.add_textbox(Inches(0.8), Inches(6.3), Inches(5), Inches(0.4))
tf = date.text_frame
p = tf.paragraphs[0]
p.text = "2026年3月"
p.font.size = Pt(14)
p.font.color.rgb = RGBColor(150, 150, 160)

# ============================================================================
# SLIDE 2: Table of Contents
# ============================================================================
slide2 = prs.slides.add_slide(prs.slide_layouts[6])

bg = slide2.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, prs.slide_height)
bg.fill.solid()
bg.fill.fore_color.rgb = CREAM
bg.line.fill.background()

add_header_accent(slide2, "目 录", "CONTENTS")
add_page_number(slide2, 2, 11)

toc_data = [
    ("01", "起源与概念界定", "Origin and Concept Definition"),
    ("02", "国际发展背景", "International Development Context"),
    ("03", "中国发展历程", "Development History in China"),
    ("04", "发展现状与规模", "Current Status and Scale"),
    ("05", "主要运作模式", "Operational Models"),
    ("06", "挑战与对策", "Challenges and Countermeasures"),
    ("07", "未来发展趋势", "Future Development Trends"),
]

y = 2.0
for num, title_cn, title_en in toc_data:
    num_box = slide2.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(y), Inches(0.8), Inches(0.6))
    num_box.fill.solid()
    num_box.fill.fore_color.rgb = NAVY
    num_box.line.fill.background()
    
    nt = slide2.shapes.add_textbox(Inches(0.8), Inches(y + 0.1), Inches(0.8), Inches(0.5))
    tf = nt.text_frame
    p = tf.paragraphs[0]
    p.text = num
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = GOLD
    p.alignment = PP_ALIGN.CENTER
    
    ct = slide2.shapes.add_textbox(Inches(1.8), Inches(y + 0.05), Inches(5), Inches(0.4))
    tf = ct.text_frame
    p = tf.paragraphs[0]
    p.text = title_cn
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = CHARCOAL
    
    et = slide2.shapes.add_textbox(Inches(1.8), Inches(y + 0.4), Inches(5), Inches(0.3))
    tf = et.text_frame
    p = tf.paragraphs[0]
    p.text = title_en
    p.font.size = Pt(11)
    p.font.color.rgb = GRAY
    
    y += 0.7

# Right decorative panel
deco_box = slide2.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(10), Inches(1.8), Inches(3), Inches(5.2))
deco_box.fill.solid()
deco_box.fill.fore_color.rgb = NAVY
deco_box.line.fill.background()

gold_line = slide2.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(10.1), Inches(2), Inches(0.06), Inches(4.8))
gold_line.fill.solid()
gold_line.fill.fore_color.rgb = GOLD
gold_line.line.fill.background()

quote = slide2.shapes.add_textbox(Inches(10.4), Inches(2.5), Inches(2.4), Inches(3))
tf = quote.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = '"日间医疗是提高医疗效率、优化资源配置的重要途径"'
p.font.size = Pt(14)
p.font.color.rgb = WHITE
p.font.italic = True

# ============================================================================
# SLIDE 3: Origin and Concept
# ============================================================================
slide3 = prs.slides.add_slide(prs.slide_layouts[6])

bg = slide3.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, prs.slide_height)
bg.fill.solid()
bg.fill.fore_color.rgb = CREAM
bg.line.fill.background()

add_header_accent(slide3, "一、起源与概念界定", "Origin and Concept Definition")
add_page_number(slide3, 3, 11)

# Definition box
def_box = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.9), Inches(12.3), Inches(1.7))
def_box.fill.solid()
def_box.fill.fore_color.rgb = WHITE
def_box.line.color.rgb = GOLD
def_box.line.width = Pt(2)

gold_accent = slide3.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.9), Inches(0.12), Inches(1.7))
gold_accent.fill.solid()
gold_accent.fill.fore_color.rgb = GOLD
gold_accent.line.fill.background()

def_title = slide3.shapes.add_textbox(Inches(0.9), Inches(2.0), Inches(11.5), Inches(0.4))
tf = def_title.text_frame
p = tf.paragraphs[0]
p.text = "【日间医疗定义】Day Care / Ambulatory Care"
p.font.size = Pt(16)
p.font.bold = True
p.font.color.rgb = NAVY

def_text = slide3.shapes.add_textbox(Inches(0.9), Inches(2.45), Inches(11.5), Inches(1))
tf = def_text.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "指患者在院时间不超过24小时的医疗服务和住院模式。患者当日入院、完成手术或操作、术后观察恢复，当日或次日出院的的一种高效医疗服务模式。"
p.font.size = Pt(15)
p.font.color.rgb = CHARCOAL

# Feature cards
features = [
    ("⏱", "住院时间短", "≤24小时"),
    ("💰", "医疗费用低", "降低成本"),
    ("🏥", "效率提升", "床位周转快"),
    ("👨‍⚕️", "流程标准化", "规范管理"),
]

x = 0.7
for icon, title, desc in features:
    box = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(3.9), Inches(2.9), Inches(1.6))
    box.fill.solid()
    box.fill.fore_color.rgb = NAVY
    box.line.fill.background()
    
    it = slide3.shapes.add_textbox(Inches(x), Inches(4.0), Inches(2.9), Inches(0.6))
    tf = it.text_frame
    p = tf.paragraphs[0]
    p.text = icon
    p.font.size = Pt(24)
    p.alignment = PP_ALIGN.CENTER
    
    tt = slide3.shapes.add_textbox(Inches(x), Inches(4.55), Inches(2.9), Inches(0.4))
    tf = tt.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    dt = slide3.shapes.add_textbox(Inches(x), Inches(4.95), Inches(2.9), Inches(0.4))
    tf = dt.text_frame
    p = tf.paragraphs[0]
    p.text = desc
    p.font.size = Pt(12)
    p.font.color.rgb = GOLD
    p.alignment = PP_ALIGN.CENTER
    
    x += 3.2

# Historical note
hist = slide3.shapes.add_textbox(Inches(0.5), Inches(5.8), Inches(12.3), Inches(0.8))
tf = hist.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "📖 历史渊源：1909年苏格兰医生James Nicoll首次提出日间手术概念 | 1995年国际日间手术协会（IAASS）成立 | 2012年中国发布首个日间手术管理规范"
p.font.size = Pt(11)
p.font.color.rgb = GRAY

# ============================================================================
# SLIDE 4: International Context
# ============================================================================
slide4 = prs.slides.add_slide(prs.slide_layouts[6])

bg = slide4.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, prs.slide_height)
bg.fill.solid()
bg.fill.fore_color.rgb = CREAM
bg.line.fill.background()

add_header_accent(slide4, "一、国际发展背景", "International Development Context")
add_page_number(slide4, 4, 11)

countries = [
    ("🇬🇧 英国", "日间手术发源地", ["1909年概念提出", "90年代全面推广", "日间手术占比超70%"], ACCENT_BLUE),
    ("🇺🇸 美国", "全球领先", ["ASC体系成熟", "日间手术占比85%+", "覆盖大部分术式"], ACCENT_GREEN),
    ("🇩🇪 德国", "标准化典范", ["1980年代开始发展", "严格日间手术标准", "高效管理系统"], RGBColor(140, 100, 180)),
    ("🇨🇳 中国", "快速发展", ["2012年规范化", "2020年扩大试点", "占比20-30%"], ACCENT_RED),
]

x = 0.5
for country, tag, items, color in countries:
    card = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(1.9), Inches(3), Inches(4.9))
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = color
    card.line.width = Pt(2)
    
    band = slide4.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(1.9), Inches(3), Inches(0.8))
    band.fill.solid()
    band.fill.fore_color.rgb = color
    band.line.fill.background()
    
    cn = slide4.shapes.add_textbox(Inches(x), Inches(2.0), Inches(3), Inches(0.6))
    tf = cn.text_frame
    p = tf.paragraphs[0]
    p.text = country
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    tg = slide4.shapes.add_textbox(Inches(x + 0.2), Inches(2.9), Inches(2.6), Inches(0.4))
    tf = tg.text_frame
    p = tf.paragraphs[0]
    p.text = "◆ " + tag
    p.font.size = Pt(13)
    p.font.color.rgb = color
    
    y = 3.4
    for item in items:
        it = slide4.shapes.add_textbox(Inches(x + 0.2), Inches(y), Inches(2.6), Inches(0.4))
        tf = it.text_frame
        p = tf.paragraphs[0]
        p.text = "• " + item
        p.font.size = Pt(12)
        p.font.color.rgb = CHARCOAL
        y += 0.5
    
    x += 3.3

# ============================================================================
# SLIDE 5: Development History
# ============================================================================
slide5 = prs.slides.add_slide(prs.slide_layouts[6])

bg = slide5.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, prs.slide_height)
bg.fill.solid()
bg.fill.fore_color.rgb = CREAM
bg.line.fill.background()

add_header_accent(slide5, "二、中国发展历程", "Development History in China")
add_page_number(slide5, 5, 11)

timeline_data = [
    ("2001-2010", "萌芽期", "• 少数大型三甲医院探索\n• 主要集中在眼科、日间化疗\n• 缺乏统一标准", NAVY),
    ("2012", "规范期", "• 原卫生部发布首个管理文件\n• 北京、上海开始试点\n• 成立日间手术合作联盟", ACCENT_BLUE),
    ("2015-2018", "推广期", "• 国家卫计委扩大试点\n• 多省市出台政策\n• 医保支付改革推进", ACCENT_GREEN),
    ("2019-至今", "快车道", "• 新技术推动（微创、麻醉）\n• 疫情加速非住院化\n• 日间病房标准化建设", GOLD),
]

tl_line = slide5.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.5), Inches(3.15), Inches(10.3), Inches(0.03))
tl_line.fill.solid()
tl_line.fill.fore_color.rgb = NAVY
tl_line.line.fill.background()

for i, (period, stage, content, color) in enumerate(timeline_data):
    dot = slide5.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.35), Inches(2.98), Inches(0.35), Inches(0.35))
    dot.fill.solid()
    dot.fill.fore_color.rgb = color
    dot.line.fill.background()
    
    y_pos = 2.2 if i % 2 == 0 else 3.5
    
    period_box = slide5.shapes.add_textbox(Inches(0.8), Inches(y_pos), Inches(1.5), Inches(0.4))
    tf = period_box.text_frame
    p = tf.paragraphs[0]
    p.text = period
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = color
    p.alignment = PP_ALIGN.CENTER
    
    stage_box = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2), Inches(y_pos + 0.4), Inches(1.8), Inches(0.5))
    stage_box.fill.solid()
    stage_box.fill.fore_color.rgb = color
    stage_box.line.fill.background()
    
    st = slide5.shapes.add_textbox(Inches(2), Inches(y_pos + 0.45), Inches(1.8), Inches(0.4))
    tf = st.text_frame
    p = tf.paragraphs[0]
    p.text = stage
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = WHITE if color != GOLD else NAVY
    p.alignment = PP_ALIGN.CENTER
    
    content_box = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4), Inches(y_pos + 0.4), Inches(8.5), Inches(1.1))
    content_box.fill.solid()
    content_box.fill.fore_color.rgb = WHITE
    content_box.line.color.rgb = color
    content_box.line.width = Pt(1)
    
    ct = slide5.shapes.add_textbox(Inches(4.2), Inches(y_pos + 0.5), Inches(8.1), Inches(0.9))
    tf = ct.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = content.replace('\n', '  |  ')
    p.font.size = Pt(11)
    p.font.color.rgb = CHARCOAL

# ============================================================================
# SLIDE 6: Current Status
# ============================================================================
slide6 = prs.slides.add_slide(prs.slide_layouts[6])

bg = slide6.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, prs.slide_height)
bg.fill.solid()
bg.fill.fore_color.rgb = CREAM
bg.line.fill.background()

add_header_accent(slide6, "三、发展现状与规模", "Current Status and Scale")
add_page_number(slide6, 6, 11)

stats = [
    ("30%+", "年均增长率", "过去5年日间手术量快速增长", ACCENT_BLUE),
    ("20-30%", "择期手术占比", "与欧美国家差距明显，提升空间大", ACCENT_GREEN),
    (">10,000", "开展机构数", "全国三级医院广泛开展日间服务", NAVY),
    ("200+", "术式种类", "覆盖普外、眼科、骨科等多个专科", GOLD),
]

y = 1.9
for i, (num, label, desc, color) in enumerate(stats):
    card = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y), Inches(6), Inches(1.2))
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = color
    card.line.width = Pt(2)
    
    band = slide6.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(y), Inches(0.12), Inches(1.2))
    band.fill.solid()
    band.fill.fore_color.rgb = color
    band.line.fill.background()
    
    num_box = slide6.shapes.add_textbox(Inches(0.8), Inches(y + 0.12), Inches(2), Inches(0.6))
    tf = num_box.text_frame
    p = tf.paragraphs[0]
    p.text = num
    p.font.size = Pt(30)
    p.font.bold = True
    p.font.color.rgb = color
    
    lb = slide6.shapes.add_textbox(Inches(0.8), Inches(y + 0.7), Inches(2), Inches(0.4))
    tf = lb.text_frame
    p = tf.paragraphs[0]
    p.text = label
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = CHARCOAL
    
    dc = slide6.shapes.add_textbox(Inches(2.9), Inches(y + 0.3), Inches(3.4), Inches(0.8))
    tf = dc.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = desc
    p.font.size = Pt(12)
    p.font.color.rgb = GRAY
    
    y += 1.3

# Right side - specialties
spec_title = slide6.shapes.add_textbox(Inches(7), Inches(1.9), Inches(5.8), Inches(0.5))
tf = spec_title.text_frame
p = tf.paragraphs[0]
p.text = "📋 覆盖主要专科"
p.font.size = Pt(16)
p.font.bold = True
p.font.color.rgb = NAVY

specialties = [
    ("🏥", "普外科", "胆囊、疝气、阑尾"),
    ("👁", "眼科", "白内障手术"),
    ("🦴", "骨科", "关节镜、日间骨折"),
    ("👃", "耳鼻喉科", "腺样体、鼻窦手术"),
    ("💊", "肿瘤科", "日间化疗"),
    ("🫁", "呼吸科", "支气管镜"),
    ("💉", "介入治疗", "血管、穿刺"),
    ("🧠", "神经外科", "介入手术"),
]

y = 2.4
for icon, dept, examples in specialties:
    row = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7), Inches(y), Inches(5.8), Inches(0.55))
    row.fill.solid()
    row.fill.fore_color.rgb = WHITE if int(y * 10) % 20 < 10 else LIGHT_GRAY
    row.line.fill.background()
    
    lb = slide6.shapes.add_textbox(Inches(7.2), Inches(y + 0.1), Inches(0.8), Inches(0.4))
    tf = lb.text_frame
    p = tf.paragraphs[0]
    p.text = icon
    p.font.size = Pt(14)
    
    dt = slide6.shapes.add_textbox(Inches(7.8), Inches(y + 0.1), Inches(1.2), Inches(0.4))
    tf = dt.text_frame
    p = tf.paragraphs[0]
    p.text = dept
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = CHARCOAL
    
    ep = slide6.shapes.add_textbox(Inches(9.0), Inches(y + 0.1), Inches(3.6), Inches(0.4))
    tf = ep.text_frame
    p = tf.paragraphs[0]
    p.text = examples
    p.font.size = Pt(11)
    p.font.color.rgb = GRAY
    
    y += 0.58

# ============================================================================
# SLIDE 7: Operational Models
# ============================================================================
slide7 = prs.slides.add_slide(prs.slide_layouts[6])

bg = slide7.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, prs.slide_height)
bg.fill.solid()
bg.fill.fore_color.rgb = CREAM
bg.line.fill.background()

add_header_accent(slide7, "四、主要运作模式", "Operational Models")
add_page_number(slide7, 7, 11)

models = [
    ("模式一", "独立日间手术中心", ["独立的日间手术中心", "专用床位和手术室", "专职医护团队", "完全按日间模式运行"], NAVY),
    ("模式二", "集中管理式", ["医院统一管理日间患者", "分散收治、统一排程", "共享手术室资源", "多学科协作(MDT)"], ACCENT_BLUE),
    ("模式三", "科室嵌入式", ["各临床科室设置日间床位", "科室内独立运作", "便于术后观察", "适合复杂病例"], ACCENT_GREEN),
]

x = 0.5
for model_num, model_name, items, color in models:
    card = slide7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(1.9), Inches(4), Inches(5.1))
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = color
    card.line.width = Pt(2)
    
    hdr = slide7.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(1.9), Inches(4), Inches(1.2))
    hdr.fill.solid()
    hdr.fill.fore_color.rgb = color
    hdr.line.fill.background()
    
    mn = slide7.shapes.add_textbox(Inches(x), Inches(2.0), Inches(4), Inches(0.5))
    tf = mn.text_frame
    p = tf.paragraphs[0]
    p.text = model_num
    p.font.size = Pt(14)
    p.font.color.rgb = GOLD
    p.alignment = PP_ALIGN.CENTER
    
    nm = slide7.shapes.add_textbox(Inches(x), Inches(2.45), Inches(4), Inches(0.6))
    tf = nm.text_frame
    p = tf.paragraphs[0]
    p.text = model_name
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    y = 3.3
    for item in items:
        blt = slide7.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x + 0.3), Inches(y + 0.12), Inches(0.12), Inches(0.12))
        blt.fill.solid()
        blt.fill.fore_color.rgb = color
        blt.line.fill.background()
        
        it = slide7.shapes.add_textbox(Inches(x + 0.55), Inches(y), Inches(3.3), Inches(0.5))
        tf = it.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = item
        p.font.size = Pt(13)
        p.font.color.rgb = CHARCOAL
        y += 0.65
    
    x += 4.3

# ============================================================================
# SLIDE 8: Challenges
# ============================================================================
slide8 = prs.slides.add_slide(prs.slide_layouts[6])

bg = slide8.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, prs.slide_height)
bg.fill.solid()
bg.fill.fore_color.rgb = CREAM
bg.line.fill.background()

add_header_accent(slide8, "四、挑战与对策", "Challenges and Countermeasures")
add_page_number(slide8, 8, 11)

challenges = [
    ("⚠️", "患者认知不足", "传统观念影响，对日间手术安全性存在担忧，术后家庭护理顾虑", "加强科普宣传，建立信任", ACCENT_RED),
    ("💳", "支付方式限制", "医保支付政策不完善，日间手术定价不合理，报销比例待提高", "推动支付方式改革", ACCENT_BLUE),
    ("🏥", "质量安全保障", "术后并发症风险，应急处理能力，随访管理难度", "完善质量保障体系", ACCENT_GREEN),
    ("🏙", "资源配置不均", "城乡差距明显，不同地区发展不平衡，基层能力不足", "促进优质资源下沉", GOLD),
]

y = 1.9
for icon, title, problem, solution, color in challenges:
    card = slide8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y), Inches(12.3), Inches(1.25))
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = color
    card.line.width = Pt(1.5)
    
    icon_box = slide8.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(y), Inches(0.8), Inches(1.25))
    icon_box.fill.solid()
    icon_box.fill.fore_color.rgb = color
    icon_box.line.fill.background()
    
    ic = slide8.shapes.add_textbox(Inches(0.5), Inches(y + 0.35), Inches(0.8), Inches(0.6))
    tf = ic.text_frame
    p = tf.paragraphs[0]
    p.text = icon
    p.font.size = Pt(22)
    p.alignment = PP_ALIGN.CENTER
    
    tt = slide8.shapes.add_textbox(Inches(1.5), Inches(y + 0.15), Inches(2.5), Inches(0.4))
    tf = tt.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = CHARCOAL
    
    pr = slide8.shapes.add_textbox(Inches(1.5), Inches(y + 0.55), Inches(6), Inches(0.6))
    tf = pr.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = problem
    p.font.size = Pt(11)
    p.font.color.rgb = GRAY
    
    arrow = slide8.shapes.add_textbox(Inches(7.8), Inches(y + 0.35), Inches(0.4), Inches(0.5))
    tf = arrow.text_frame
    p = tf.paragraphs[0]
    p.text = "→"
    p.font.size = Pt(20)
    p.font.color.rgb = color
    
    sol_box = slide8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.3), Inches(y + 0.25), Inches(4.3), Inches(0.75))
    sol_box.fill.solid()
    sol_box.fill.fore_color.rgb = color
    sol_box.line.fill.background()
    
    sl = slide8.shapes.add_textbox(Inches(8.3), Inches(y + 0.4), Inches(4.3), Inches(0.5))
    tf